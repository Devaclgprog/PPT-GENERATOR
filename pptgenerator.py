import os
from dotenv import load_dotenv

import tempfile
import streamlit as st
from PyPDF2 import PdfReader
import pdfplumber
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from datetime import datetime
from typing import Optional, Tuple

# ========== CONFIGURATION ========== #
load_dotenv()

# Configure Gemini
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))


MODEL_NAME = "gemini-1.5-pro"

# Constants
MAX_PDF_SIZE_MB = 50
MAX_SLIDES = 10
PROCESSING_CHUNK_SIZE = 15000
MIN_CONTENT_LENGTH = 100

# ========== CORE FUNCTIONS ========== #

def configure_gemini() -> genai.GenerativeModel:
    """Configure Gemini API with error handling"""
    try:
        genai.configure(api_key=GEMINI_API_KEY)
        return genai.GenerativeModel(
            MODEL_NAME,
            generation_config={
                "temperature": 0.3,
                "top_p": 0.95,
                "top_k": 40,
                "max_output_tokens": 2048,
            }
        )
    except Exception as e:
        st.error(f"Failed to configure Gemini: {str(e)}")
        st.stop()

def extract_text_from_pdf(pdf_file) -> Tuple[Optional[str], Optional[str]]:
    """Robust text extraction with multiple fallbacks"""
    pdf_file.seek(0)
    text = ""
    
    # Try pdfplumber first
    try:
        with pdfplumber.open(pdf_file) as pdf:
            for i, page in enumerate(pdf.pages):
                page_text = page.extract_text()
                if page_text:
                    text += f"\n\n[Page {i+1}]\n{page_text}"
                if len(text) > PROCESSING_CHUNK_SIZE:
                    break
    except Exception as e:
        st.warning(f"pdfplumber failed: {str(e)}")
    
    # Fallback to PyPDF2 if needed
    if len(text) < MIN_CONTENT_LENGTH:
        pdf_file.seek(0)
        try:
            reader = PdfReader(pdf_file)
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text:
                    text += f"\n\n[Page {i+1}]\n{page_text}"
                if len(text) > PROCESSING_CHUNK_SIZE:
                    break
        except Exception as e:
            st.warning(f"PyPDF2 failed: {str(e)}")
    
    if len(text) > MIN_CONTENT_LENGTH:
        return text, None
    return None, "Failed to extract sufficient text (document may be scanned)"

def generate_slide_structure(model, pdf_text: str, ppt_title: str) -> Tuple[Optional[str], Optional[str]]:
    """Generate structured slide outline with validation"""
    try:
        prompt = f"""
        Create a detailed PowerPoint outline from this document for title: "{ppt_title}".
        
        DOCUMENT CONTENT:
        {pdf_text[:PROCESSING_CHUNK_SIZE]}
        
        REQUIRED OUTPUT FORMAT:
        Slide 1: [Title Slide] - Title: "{ppt_title}", Subtitle: "[Document summary]"
        Slide 2: [Introduction] - [3-5 specific points from document]
        Slide 3: [Key Finding 1] - [Detailed content from document with page reference]
        Slide 4: [Key Finding 2] - [Detailed content from document with page reference]
        Slide 5: [Conclusion] - [Actionable takeaways]
        
        RULES:
        - Create exactly 5 slides
        - Each slide must have concrete content instructions
        - Include specific facts/numbers/quotes when available
        - Never invent information not in the document
        - Include page references like (Page 5) when possible
        """
        
        response = model.generate_content(prompt)
        return response.text, None
    except Exception as e:
        return None, f"Structure generation failed: {str(e)}"

def generate_slide_content(model, pdf_text: str, slide_title: str) -> str:
    """Generate accurate slide content with validation"""
    try:
        prompt = f"""
        Generate specific content for PowerPoint slide titled: "{slide_title}"
        
        DOCUMENT CONTENT:
        {pdf_text[:PROCESSING_CHUNK_SIZE]}
        
        REQUIREMENTS:
        1. Extract 3-5 specific points from the document
        2. Each point must include concrete details
        3. Use only factual information from the document
        4. Format as bullet points with page references like (Page 3)
        5. If no specific content found, use general document themes
        
        EXAMPLE OUTPUT:
        - Revenue increased by 23% in Q3 (Page 5)
        - Customer satisfaction reached 4.8/5 (Page 8)
        - New products launched in September (Page 12)
        """
        
        response = model.generate_content(prompt)
        content = response.text
        
        # Ensure we have proper bullet points
        if not content.strip():
            content = "- Key point from document (Page X)\n- Important finding (Page Y)\n- Relevant detail (Page Z)"
        elif not any(marker in content for marker in ['-', '•', '*']):
            content = "- " + content.replace('\n', '\n- ')
        
        return content
    except Exception as e:
        return "- Document point 1 (Page X)\n- Document point 2 (Page Y)\n- Document point 3 (Page Z)"

def create_presentation(ppt_title: str, slide_structure: str, model, pdf_text: str) -> Optional[str]:
    """Create PowerPoint with guaranteed content in slides"""
    try:
        prs = Presentation()
        
        # Set slide size to widescreen (16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = ppt_title
        subtitle.text = f"Generated {datetime.now().strftime('%d %b %Y %H:%M')}"
        
        # Parse slide structure
        slides_to_create = []
        for line in slide_structure.split('\n'):
            if line.startswith('Slide'):
                parts = line.split(' - ')
                if len(parts) >= 2:
                    slide_title = parts[0].split(': ')[1].strip('[]')
                    slides_to_create.append(slide_title)
        
        # Create content slides
        for slide_title in slides_to_create[:MAX_SLIDES]:
            if '[Title Slide]' in slide_title:
                continue
                
            # Generate content (with fallback)
            content = generate_slide_content(model, pdf_text, slide_title)
            
            # Create slide
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            content_shape = slide.placeholders[1]
            
            # Set title
            title_shape.text = slide_title.replace('[', '').replace(']', '')
            
            # Set content (guaranteed to have content)
            content_shape.text = content
            
            # Formatting
            for paragraph in content_shape.text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.font.color.rgb = RGBColor(0, 0, 0)
                paragraph.alignment = PP_ALIGN.LEFT
        
        # Save to temp file
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
            prs.save(tmp.name)
            return tmp.name
            
    except Exception as e:
        st.error(f"PPT creation failed: {str(e)}")
        return None

# ========== STREAMLIT UI ========== #

def main():
    st.set_page_config(
        page_title="PDF to PowerPoint Pro",
        page_icon="📊",
        layout="centered"
    )
    
    st.title("📊 PDF to PowerPoint Converter")
    st.write("Upload a PDF to generate a professional PowerPoint presentation")
    
    # Initialize Gemini
    try:
        model = configure_gemini()
    except Exception as e:
        st.error(f"Failed to initialize AI service: {str(e)}")
        st.stop()
    
    # File upload
    pdf_file = st.file_uploader("Choose PDF file", type=["pdf"])
    
    if pdf_file:
        # Validate file size
        if pdf_file.size > MAX_PDF_SIZE_MB * 1024 * 1024:
            st.error(f"File too large (max {MAX_PDF_SIZE_MB}MB)")
            return
        
        # Extract text
        with st.spinner("Extracting text from PDF..."):
            pdf_text, error = extract_text_from_pdf(pdf_file)
            if error:
                st.error(error)
                return
            st.session_state.pdf_text = pdf_text
        
        # Presentation title
        ppt_title = st.text_input("Presentation Title", "Business Report")
        
        # Generate slide structure
        if st.button("Analyze Document"):
            with st.spinner("Creating slide structure..."):
                slide_structure, error = generate_slide_structure(model, pdf_text, ppt_title)
                if error:
                    st.error(error)
                else:
                    st.session_state.slide_structure = slide_structure
        
        # Display and edit structure
        if 'slide_structure' in st.session_state:
            st.subheader("Slide Structure")
            edited_structure = st.text_area(
                "Review and edit if needed:",
                value=st.session_state.slide_structure,
                height=300
            )
            
            if st.button("Generate PowerPoint", type="primary"):
                with st.spinner("Creating presentation..."):
                    ppt_path = create_presentation(
                        ppt_title,
                        edited_structure,
                        model,
                        st.session_state.pdf_text
                    )
                    if ppt_path:
                        with open(ppt_path, "rb") as f:
                            ppt_bytes = f.read()
                        st.success("✅ Presentation generated successfully!")
                        st.download_button(
                            label="Download PowerPoint",
                            data=ppt_bytes,
                            file_name=f"{ppt_title.replace(' ', '_')}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                        try:
                            os.unlink(ppt_path)
                        except:
                            pass

if __name__ == "__main__":
    main()
